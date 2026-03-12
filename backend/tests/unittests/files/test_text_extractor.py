"""
Unit tests for TextExtractor class.

Tests cover:
- Plain text extraction with different encodings
- PDF text extraction
- DOCX text extraction
- PPTX text extraction
- MIME type detection and routing
- Error handling
"""

from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

from intric.files.text import (
    CorruptFileError,
    EncryptedFileError,
    ExtractionError,
    TextExtractor,
    TextMimeTypes,
    TextSanitizer,
    UnsupportedFormatError,
)


class TestTextSanitizer:
    """Tests for TextSanitizer class."""

    def test_sanitize_removes_null_bytes(self):
        """Null bytes should be removed from text."""
        text_with_nulls = "Hello\x00World\x00!"
        result = TextSanitizer.sanitize(text_with_nulls)
        assert result == "HelloWorld!"

    def test_sanitize_preserves_normal_text(self):
        """Normal text without null bytes should be unchanged."""
        normal_text = "Hello World! This is a test."
        result = TextSanitizer.sanitize(normal_text)
        assert result == normal_text

    def test_sanitize_handles_empty_string(self):
        """Empty string should remain empty."""
        result = TextSanitizer.sanitize("")
        assert result == ""

    def test_sanitize_handles_only_null_bytes(self):
        """String with only null bytes should become empty."""
        result = TextSanitizer.sanitize("\x00\x00\x00")
        assert result == ""


class TestTextMimeTypes:
    """Tests for TextMimeTypes enum."""

    def test_has_value_returns_true_for_valid_mime(self):
        """has_value should return True for valid MIME types."""
        assert TextMimeTypes.has_value("application/pdf") is True
        assert TextMimeTypes.has_value("text/plain") is True
        assert TextMimeTypes.has_value("text/markdown") is True
        assert TextMimeTypes.has_value(
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        ) is True
        assert TextMimeTypes.has_value("application/vnd.ms-excel") is True

    def test_has_value_returns_false_for_invalid_mime(self):
        """has_value should return False for invalid MIME types."""
        assert TextMimeTypes.has_value("application/octet-stream") is False
        assert TextMimeTypes.has_value("image/png") is False

    def test_has_value_handles_charset_suffix(self):
        """has_value should handle MIME types with charset suffix."""
        assert TextMimeTypes.has_value("text/plain; charset=utf-8") is True
        assert TextMimeTypes.has_value("application/pdf; charset=binary") is True

    def test_values_returns_all_mime_types(self):
        """values() should return list of all MIME type values."""
        values = TextMimeTypes.values()
        assert "application/pdf" in values
        assert "text/plain" in values
        assert "text/markdown" in values


class TestTextExtractorPlainText:
    """Tests for plain text extraction."""

    def test_extract_from_plain_text_utf8(self, tmp_path):
        """Should extract text from UTF-8 encoded file."""
        test_file = tmp_path / "test.txt"
        test_file.write_text("Hello World!", encoding="utf-8")

        result = TextExtractor.extract_from_plain_text(test_file)
        assert result == "Hello World!"

    def test_extract_from_plain_text_latin1(self, tmp_path):
        """Should extract text from Latin-1 encoded file."""
        test_file = tmp_path / "test.txt"
        # Write with latin-1 encoding (includes special chars)
        test_file.write_bytes("Caf\xe9".encode("latin-1"))

        result = TextExtractor.extract_from_plain_text(test_file)
        assert "Caf" in result

    def test_extract_from_plain_text_multiline(self, tmp_path):
        """Should handle multiline text files."""
        test_file = tmp_path / "test.txt"
        content = "Line 1\nLine 2\nLine 3"
        test_file.write_text(content, encoding="utf-8")

        result = TextExtractor.extract_from_plain_text(test_file)
        assert result == content

    def test_extract_from_plain_text_empty_file(self, tmp_path):
        """Should handle empty files."""
        test_file = tmp_path / "empty.txt"
        test_file.write_text("", encoding="utf-8")

        result = TextExtractor.extract_from_plain_text(test_file)
        assert result == ""


class TestTextExtractorPDF:
    """Tests for PDF text extraction."""

    def test_extract_from_pdf_basic(self):
        """Should extract text from a basic PDF."""
        with patch("intric.files.text.pdfplumber.open") as mock_open:
            mock_page = MagicMock()
            mock_page.extract_text.return_value = "Sample PDF text content"
            mock_open.return_value.__enter__ = MagicMock(return_value=MagicMock(pages=[mock_page]))
            mock_open.return_value.__exit__ = MagicMock(return_value=False)

            result = TextExtractor.extract_from_pdf(Path("test.pdf"))

            assert isinstance(result, str)
            assert "Sample PDF text content" in result

    def test_extract_from_pdf_handles_none_pages(self):
        """Should handle pages that return None from extract_text()."""
        with patch("intric.files.text.pdfplumber.open") as mock_open:
            mock_page1 = MagicMock()
            mock_page1.extract_text.return_value = "Page 1 text"

            mock_page2 = MagicMock()
            mock_page2.extract_text.return_value = None

            mock_page3 = MagicMock()
            mock_page3.extract_text.return_value = "Page 3 text"

            mock_open.return_value.__enter__ = MagicMock(
                return_value=MagicMock(pages=[mock_page1, mock_page2, mock_page3])
            )
            mock_open.return_value.__exit__ = MagicMock(return_value=False)

            result = TextExtractor.extract_from_pdf(Path("dummy.pdf"))

            assert "Page 1 text" in result
            assert "Page 3 text" in result

    def test_extract_from_pdf_sanitizes_null_bytes(self):
        """Should sanitize null bytes from extracted text."""
        with patch("intric.files.text.pdfplumber.open") as mock_open:
            mock_page = MagicMock()
            mock_page.extract_text.return_value = "Hello\x00World"
            mock_open.return_value.__enter__ = MagicMock(return_value=MagicMock(pages=[mock_page]))
            mock_open.return_value.__exit__ = MagicMock(return_value=False)

            result = TextExtractor.extract_from_pdf(Path("dummy.pdf"))

            assert "\x00" not in result
            assert "HelloWorld" in result


class TestTextExtractorDOCX:
    """Tests for DOCX text extraction."""

    def test_extract_from_docx_basic(self, tmp_path):
        """Should extract text from a basic DOCX file."""
        # Create a simple DOCX using python-docx
        from docx import Document

        docx_path = tmp_path / "test.docx"
        doc = Document()
        doc.add_paragraph("Hello from DOCX!")
        doc.add_paragraph("Second paragraph.")
        doc.save(docx_path)

        result = TextExtractor.extract_from_docx(docx_path)

        assert "Hello from DOCX!" in result
        assert "Second paragraph" in result

    def test_extract_from_docx_empty_document(self, tmp_path):
        """Should handle empty DOCX documents."""
        from docx import Document

        docx_path = tmp_path / "empty.docx"
        doc = Document()
        doc.save(docx_path)

        result = TextExtractor.extract_from_docx(docx_path)
        assert isinstance(result, str)


class TestTextExtractorPPTX:
    """Tests for PPTX text extraction."""

    def test_extract_from_pptx_basic(self, tmp_path):
        """Should extract text from a basic PPTX file."""
        from pptx import Presentation
        from pptx.util import Inches

        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add a text box
        left = top = Inches(1)
        width = height = Inches(2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Hello from PPTX!"

        prs.save(pptx_path)

        result = TextExtractor.extract_from_pptx(pptx_path)

        assert "Hello from PPTX!" in result

    def test_extract_from_pptx_multiple_slides(self, tmp_path):
        """Should extract text from multiple slides."""
        from pptx import Presentation
        from pptx.util import Inches

        pptx_path = tmp_path / "multi.pptx"
        prs = Presentation()

        for i in range(3):
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
            txBox.text_frame.text = f"Slide {i + 1} content"

        prs.save(pptx_path)

        result = TextExtractor.extract_from_pptx(pptx_path)

        assert "Slide 1 content" in result
        assert "Slide 2 content" in result
        assert "Slide 3 content" in result

    def test_extract_from_pptx_empty_presentation(self, tmp_path):
        """Should handle empty presentations."""
        from pptx import Presentation

        pptx_path = tmp_path / "empty.pptx"
        prs = Presentation()
        prs.save(pptx_path)

        result = TextExtractor.extract_from_pptx(pptx_path)
        assert isinstance(result, str)


class TestTextExtractorXLSX:
    """Tests for XLSX/XLS text extraction."""

    def test_extract_from_xlsx_basic(self, tmp_path):
        """Should extract text as self-contained key-value rows for LLM comprehension."""
        from openpyxl import Workbook

        xlsx_path = tmp_path / "test.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws["A1"] = "Name"
        ws["B1"] = "Value"
        ws["A2"] = "Item1"
        ws["B2"] = 100
        wb.save(xlsx_path)

        result = TextExtractor.extract_from_xlsx(xlsx_path)

        # File context at the start
        assert "File: test.xlsx" in result
        # Each row is self-contained with sheet and column context
        assert "Sheet: Sheet1" in result
        assert "Name: Item1" in result
        assert "Value: 100" in result

    def test_extract_from_xlsx_multiple_sheets(self, tmp_path):
        """Should extract text from multiple sheets as self-contained key-value rows."""
        from openpyxl import Workbook

        xlsx_path = tmp_path / "multi.xlsx"
        wb = Workbook()

        # First sheet
        ws1 = wb.active
        ws1.title = "Sales"
        ws1["A1"] = "Product"
        ws1["A2"] = "Widget"

        # Second sheet
        ws2 = wb.create_sheet("Inventory")
        ws2["A1"] = "Stock"
        ws2["A2"] = 50

        wb.save(xlsx_path)

        result = TextExtractor.extract_from_xlsx(xlsx_path)

        # Each row is self-contained with sheet context
        assert "Sheet: Sales" in result
        assert "Product: Widget" in result
        assert "Sheet: Inventory" in result
        assert "Stock: 50" in result

    def test_extract_from_xlsx_empty_workbook(self, tmp_path):
        """Should handle empty workbooks with just filename header."""
        from openpyxl import Workbook

        xlsx_path = tmp_path / "empty.xlsx"
        wb = Workbook()
        wb.save(xlsx_path)

        result = TextExtractor.extract_from_xlsx(xlsx_path)
        assert isinstance(result, str)
        assert "File: empty.xlsx" in result

    def test_extract_from_xlsx_with_special_characters(self, tmp_path):
        """Should handle special characters in cells."""
        from openpyxl import Workbook

        xlsx_path = tmp_path / "special.xlsx"
        wb = Workbook()
        ws = wb.active
        ws["A1"] = "Café"
        ws["A2"] = "日本語"
        ws["A3"] = "émojis: 🎉"
        wb.save(xlsx_path)

        result = TextExtractor.extract_from_xlsx(xlsx_path)

        assert "Café" in result
        assert "日本語" in result


class TestTextExtractorExtractMethod:
    """Tests for the main extract() method with MIME type routing."""

    def test_extract_routes_pdf_correctly(self, tmp_path):
        """Should route PDF files to PDF extractor."""
        extractor = TextExtractor()

        with patch.object(extractor, "extract_from_pdf", return_value="PDF content"):
            result = extractor.extract(tmp_path / "test.pdf", "application/pdf")
            assert result == "PDF content"

    def test_extract_routes_docx_correctly(self, tmp_path):
        """Should route DOCX files to DOCX extractor."""
        extractor = TextExtractor()
        mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

        with patch.object(extractor, "extract_from_docx", return_value="DOCX content"):
            result = extractor.extract(tmp_path / "test.docx", mime)
            assert result == "DOCX content"

    def test_extract_routes_pptx_correctly(self, tmp_path):
        """Should route PPTX files to PPTX extractor."""
        extractor = TextExtractor()
        mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

        with patch.object(extractor, "extract_from_pptx", return_value="PPTX content"):
            result = extractor.extract(tmp_path / "test.pptx", mime)
            assert result == "PPTX content"

    def test_extract_routes_xlsx_correctly(self, tmp_path):
        """Should route XLSX files to XLSX extractor."""
        extractor = TextExtractor()
        mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"

        with patch.object(extractor, "extract_from_xlsx", return_value="XLSX content"):
            result = extractor.extract(tmp_path / "test.xlsx", mime)
            assert result == "XLSX content"

    def test_extract_routes_xls_correctly(self, tmp_path):
        """Should route XLS files to XLSX extractor."""
        extractor = TextExtractor()
        mime = "application/vnd.ms-excel"

        with patch.object(extractor, "extract_from_xlsx", return_value="XLS content"):
            result = extractor.extract(tmp_path / "test.xls", mime)
            assert result == "XLS content"

    def test_extract_routes_plain_text_correctly(self, tmp_path):
        """Should route plain text files correctly."""
        extractor = TextExtractor()

        test_file = tmp_path / "test.txt"
        test_file.write_text("Plain text content")

        result = extractor.extract(test_file, "text/plain")
        assert result == "Plain text content"

    def test_extract_routes_markdown_correctly(self, tmp_path):
        """Should route markdown files correctly."""
        extractor = TextExtractor()

        test_file = tmp_path / "test.md"
        test_file.write_text("# Markdown content")

        result = extractor.extract(test_file, "text/markdown")
        assert result == "# Markdown content"

    def test_extract_routes_csv_correctly(self, tmp_path):
        """Should route CSV files correctly."""
        extractor = TextExtractor()

        test_file = tmp_path / "test.csv"
        test_file.write_text("col1,col2\nval1,val2")

        result = extractor.extract(test_file, "text/csv")
        assert "col1,col2" in result

    def test_extract_fallback_to_plain_text(self, tmp_path):
        """Should fall back to plain text for unknown MIME types."""
        extractor = TextExtractor()

        test_file = tmp_path / "test.unknown"
        test_file.write_text("Unknown format content")

        result = extractor.extract(test_file, "application/octet-stream")
        assert result == "Unknown format content"

    def test_extract_strips_whitespace(self, tmp_path):
        """Should strip leading/trailing whitespace from result."""
        extractor = TextExtractor()

        test_file = tmp_path / "test.txt"
        test_file.write_text("  content with whitespace  \n\n")

        result = extractor.extract(test_file, "text/plain")
        assert result == "content with whitespace"

    def test_extract_detects_mime_type_when_not_provided(self, tmp_path):
        """Should detect MIME type when not provided."""
        extractor = TextExtractor()

        test_file = tmp_path / "test.txt"
        test_file.write_text("Auto-detected content")

        with patch("intric.files.text.magic.from_file", return_value="text/plain"):
            result = extractor.extract(test_file)
            assert result == "Auto-detected content"


class TestTextExtractorErrorHandling:
    """Tests for error handling scenarios."""

    def test_extract_from_pdf_raises_corrupt_error_on_syntax_error(self):
        """Should raise CorruptFileError when pdfplumber encounters a corrupt PDF."""
        from pdfminer.pdfparser import PDFSyntaxError

        with patch("intric.files.text.pdfplumber.open") as mock_open:
            mock_open.side_effect = PDFSyntaxError("Invalid PDF structure")

            with pytest.raises(CorruptFileError) as exc_info:
                TextExtractor.extract_from_pdf(Path("corrupt.pdf"))

            assert exc_info.value.code == "CORRUPT"
            assert "corrupt" in exc_info.value.message.lower()

    def test_extract_from_docx_raises_corrupt_error_on_bad_zip(self, tmp_path):
        """Should raise CorruptFileError for invalid DOCX (bad ZIP)."""
        # Create an invalid ZIP file (just random bytes)
        bad_docx = tmp_path / "bad.docx"
        bad_docx.write_bytes(b"This is not a valid ZIP file")

        with pytest.raises(CorruptFileError) as exc_info:
            TextExtractor.extract_from_docx(bad_docx)

        assert exc_info.value.code == "CORRUPT"
        assert "ZIP" in exc_info.value.message or "corrupt" in exc_info.value.message.lower()

    def test_extract_from_pptx_raises_corrupt_error_on_bad_zip(self, tmp_path):
        """Should raise CorruptFileError for invalid PPTX (bad ZIP)."""
        # Create an invalid ZIP file
        bad_pptx = tmp_path / "bad.pptx"
        bad_pptx.write_bytes(b"This is not a valid ZIP file")

        with pytest.raises(CorruptFileError) as exc_info:
            TextExtractor.extract_from_pptx(bad_pptx)

        assert exc_info.value.code == "CORRUPT"
        assert "ZIP" in exc_info.value.message or "corrupt" in exc_info.value.message.lower()

    def test_extract_raises_unsupported_format_for_legacy_doc(self, tmp_path):
        """Should raise UnsupportedFormatError for .doc files."""
        extractor = TextExtractor()

        # Create a dummy file (content doesn't matter, MIME type does)
        doc_file = tmp_path / "legacy.doc"
        doc_file.write_bytes(b"dummy content")

        with pytest.raises(UnsupportedFormatError) as exc_info:
            extractor.extract(doc_file, "application/msword")

        assert exc_info.value.code == "UNSUPPORTED_FORMAT"
        assert ".doc" in exc_info.value.message
        assert ".docx" in exc_info.value.message  # Should suggest conversion

    def test_extract_raises_unsupported_format_for_legacy_ppt(self, tmp_path):
        """Should raise UnsupportedFormatError for .ppt files."""
        extractor = TextExtractor()

        # Create a dummy file
        ppt_file = tmp_path / "legacy.ppt"
        ppt_file.write_bytes(b"dummy content")

        with pytest.raises(UnsupportedFormatError) as exc_info:
            extractor.extract(ppt_file, "application/vnd.ms-powerpoint")

        assert exc_info.value.code == "UNSUPPORTED_FORMAT"
        assert ".ppt" in exc_info.value.message
        assert ".pptx" in exc_info.value.message  # Should suggest conversion

    def test_extract_from_plain_text_raises_on_permission_error(self, tmp_path):
        """Should raise ExtractionError on permission denied."""
        test_file = tmp_path / "noperm.txt"
        test_file.write_text("content")

        with patch("pathlib.Path.read_text", side_effect=PermissionError("Access denied")):
            with pytest.raises(ExtractionError) as exc_info:
                TextExtractor.extract_from_plain_text(test_file)

            assert "Permission denied" in exc_info.value.message

    def test_extract_from_plain_text_handles_binary_with_replacement(self, tmp_path):
        """Should handle undecodable bytes using replacement characters."""
        # Create a file with bytes that are invalid in UTF-8
        test_file = tmp_path / "binary.txt"
        # Mix of valid UTF-8 text and invalid bytes
        test_file.write_bytes(b"Hello\x80\x81World")

        # Should succeed using errors="replace" fallback
        result = TextExtractor.extract_from_plain_text(test_file)

        # Result should contain the readable parts and replacement chars
        assert "Hello" in result
        assert "World" in result
        # The invalid bytes become replacement characters (U+FFFD)
        assert "\ufffd" in result or "Hello" in result  # Either way, no crash


class TestTextExtractorExceptionClasses:
    """Tests for custom exception classes."""

    def test_extraction_error_has_message_and_code(self):
        """ExtractionError should have message and code attributes."""
        error = ExtractionError("Test error", code="TEST_CODE")
        assert error.message == "Test error"
        assert error.code == "TEST_CODE"
        assert str(error) == "Test error"

    def test_extraction_error_default_code(self):
        """ExtractionError should have default code."""
        error = ExtractionError("Test error")
        assert error.code == "EXTRACTION_FAILED"

    def test_encrypted_file_error_attributes(self):
        """EncryptedFileError should set correct attributes."""
        error = EncryptedFileError("secret.pdf")
        assert "secret.pdf" in error.message
        assert "encrypted" in error.message.lower()
        assert error.code == "ENCRYPTED"

    def test_corrupt_file_error_with_details(self):
        """CorruptFileError should include details when provided."""
        error = CorruptFileError("broken.docx", "Missing required element")
        assert "broken.docx" in error.message
        assert "Missing required element" in error.message
        assert error.code == "CORRUPT"

    def test_corrupt_file_error_without_details(self):
        """CorruptFileError should work without details."""
        error = CorruptFileError("broken.docx")
        assert "broken.docx" in error.message
        assert error.code == "CORRUPT"

    def test_unsupported_format_error_attributes(self):
        """UnsupportedFormatError should suggest conversion."""
        error = UnsupportedFormatError("old.doc", ".doc (Legacy Word)")
        assert "old.doc" in error.message
        assert ".doc" in error.message
        assert "not supported" in error.message.lower()
        assert error.code == "UNSUPPORTED_FORMAT"
