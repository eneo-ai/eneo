import logging
import zipfile
from enum import Enum
from pathlib import Path

import magic
import pptx
from docx2python import docx2python
import pdfplumber
from pdfminer.pdfparser import PDFSyntaxError

logger = logging.getLogger(__name__)


# =============================================================================
# Custom Exceptions
# =============================================================================


class ExtractionError(Exception):
    """Base exception for text extraction failures."""

    def __init__(self, message: str, code: str = "EXTRACTION_FAILED"):
        self.message = message
        self.code = code
        super().__init__(self.message)


class EncryptedFileError(ExtractionError):
    """Raised for password-protected files."""

    def __init__(self, filename: str):
        super().__init__(
            f"File '{filename}' is encrypted/password-protected", "ENCRYPTED"
        )


class CorruptFileError(ExtractionError):
    """Raised for corrupted or malformed files."""

    def __init__(self, filename: str, details: str = ""):
        message = f"File '{filename}' is corrupted or malformed"
        if details:
            message = f"{message}. {details}"
        super().__init__(message, "CORRUPT")


class UnsupportedFormatError(ExtractionError):
    """Raised for unsupported file formats."""

    def __init__(self, filename: str, format_type: str):
        super().__init__(
            f"Format '{format_type}' not supported. "
            f"Please convert '{filename}' to a supported format.",
            "UNSUPPORTED_FORMAT",
        )


# =============================================================================
# MIME Types
# =============================================================================


class MimeTypesBase(str, Enum):
    @classmethod
    def has_value(cls, value) -> bool:
        base_value = value.split(";")[0].strip()
        return any(base_value == item.value for item in cls)

    @classmethod
    def values(cls):
        return [item.value for item in cls]


class TextMimeTypes(MimeTypesBase):
    # Supported formats
    MD = "text/markdown"
    TXT = "text/plain"
    PDF = "application/pdf"
    DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    TEXT_CSV = "text/csv"
    APP_CSV = "application/csv"
    PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    XLS = "application/vnd.ms-excel"

    # Legacy formats (for detection/rejection only)
    DOC = "application/msword"
    PPT = "application/vnd.ms-powerpoint"


# =============================================================================
# Text Processing
# =============================================================================


class TextSanitizer:
    @staticmethod
    def sanitize(text: str) -> str:
        text = text.replace("\x00", "")
        return text


class TextExtractor:
    @staticmethod
    def extract_from_plain_text(filepath: Path, filename: str | None = None) -> str:
        display_name = filename or filepath.name
        # Try UTF-8 first, then cp1252 (Windows), then UTF-8 with replacement
        # Avoid latin-1 as it accepts any byte sequence including binary garbage
        try:
            return filepath.read_text("utf-8")
        except UnicodeDecodeError:
            pass
        except PermissionError as e:
            raise ExtractionError(f"Permission denied reading '{display_name}': {e}")
        except OSError as e:
            raise ExtractionError(f"Error reading '{display_name}': {e}")

        try:
            return filepath.read_text("cp1252")
        except UnicodeDecodeError:
            pass

        # Final fallback: UTF-8 with replacement characters for undecodable bytes
        try:
            return filepath.read_text("utf-8", errors="replace")
        except (PermissionError, OSError) as e:
            raise ExtractionError(f"Error reading '{display_name}': {e}")

    @staticmethod
    def extract_from_pdf(filepath: Path, filename: str | None = None) -> str:
        display_name = filename or filepath.name
        try:
            with pdfplumber.open(filepath) as pdf:
                extracted_text = " ".join(
                    page.extract_text() or "" for page in pdf.pages
                )

            # Warn if no text extracted (likely image-only/scanned PDF)
            sanitized = TextSanitizer.sanitize(extracted_text)
            if not sanitized.strip():
                logger.warning(
                    f"No text extracted from PDF '{display_name}' - "
                    "file may be image-only or scanned"
                )

            return sanitized

        except PDFSyntaxError as e:
            logger.warning(f"PDF read error for {display_name}: {e}")
            raise CorruptFileError(display_name, str(e))
        except Exception as e:
            logger.error(f"Unexpected PDF extraction error for {display_name}: {e}")
            raise ExtractionError(f"PDF extraction failed for '{display_name}': {str(e)}")

    @staticmethod
    def extract_from_docx(filepath: Path, filename: str | None = None) -> str:
        display_name = filename or filepath.name
        try:
            with docx2python(filepath) as docx_content:
                return docx_content.text
        except zipfile.BadZipFile:
            raise CorruptFileError(
                display_name,
                "Invalid ZIP structure - file may be corrupted or in legacy .doc format",
            )
        except KeyError as e:
            raise CorruptFileError(
                display_name, f"Missing required document component: {e}"
            )
        except Exception as e:
            logger.error(f"Unexpected DOCX extraction error for {display_name}: {e}")
            raise ExtractionError(f"DOCX extraction failed for '{display_name}': {str(e)}")

    @staticmethod
    def extract_from_xlsx(filepath: Path, filename: str | None = None) -> str:
        import pandas as pd

        display_name = filename or filepath.name
        try:
            xls = pd.ExcelFile(filepath, engine="calamine")
            parts = []

            # Global file context (helpful for first chunk and direct chat)
            parts.append(f"File: {display_name}")

            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name, engine="calamine")

                if df.empty:
                    continue

                # Handle merged cells - forward fill values
                df = df.ffill()

                # Clean column names (ensure strings, remove newlines)
                df.columns = df.columns.astype(str).str.replace("\n", " ")

                # Serialize each row as self-contained key-value pairs
                # This ensures every chunk has full context, even if split
                def serialize_row(row):
                    # Filter out NaN to save tokens
                    pairs = [
                        f"{col}: {val}" for col, val in row.items() if pd.notna(val)
                    ]
                    return f"Sheet: {sheet_name} | " + " | ".join(pairs)

                sheet_text = df.apply(serialize_row, axis=1).str.cat(sep="\n")
                parts.append(sheet_text)

            return "\n\n".join(parts)
        except ValueError as e:
            raise CorruptFileError(display_name, f"Cannot parse Excel format: {e}")
        except Exception as e:
            logger.error(f"Unexpected Excel extraction error for {display_name}: {e}")
            raise ExtractionError(f"Excel extraction failed for '{display_name}': {str(e)}")

    @staticmethod
    def extract_from_pptx(filepath: Path, filename: str | None = None) -> str:
        display_name = filename or filepath.name
        try:
            # Extract text from pptx using python-pptx
            # Use list join instead of string concatenation for O(n) vs O(n^2) complexity
            presentation = pptx.Presentation(filepath)
            parts = []
            for slide in presentation.slides:
                slide_parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        # Collect all text from runs in this shape
                        shape_text = " ".join(
                            run.text
                            for para in shape.text_frame.paragraphs
                            for run in para.runs
                            if run.text
                        )
                        if shape_text.strip():
                            slide_parts.append(shape_text)
                if slide_parts:
                    parts.append(" ".join(slide_parts))
            return "\n".join(parts)
        except zipfile.BadZipFile:
            raise CorruptFileError(
                display_name,
                "Invalid ZIP structure - file may be corrupted or in legacy .ppt format",
            )
        except KeyError as e:
            raise CorruptFileError(
                display_name, f"Missing required presentation component: {e}"
            )
        except Exception as e:
            logger.error(f"Unexpected PPTX extraction error for {display_name}: {e}")
            raise ExtractionError(f"PPTX extraction failed for '{display_name}': {str(e)}")

    def extract(
        self, filepath: Path, mimetype: str | None = None, filename: str | None = None
    ) -> str:
        mimetype = mimetype or magic.from_file(filepath, mime=True)
        # Use original filename for error messages, fallback to temp filepath
        display_name = filename or filepath.name

        # Reject legacy formats early with helpful message
        if mimetype == TextMimeTypes.DOC.value:
            raise UnsupportedFormatError(
                display_name,
                ".doc (Legacy Word) - please save as .docx",
            )
        if mimetype == TextMimeTypes.PPT.value:
            raise UnsupportedFormatError(
                display_name,
                ".ppt (Legacy PowerPoint) - please save as .pptx",
            )

        match mimetype:
            case (
                TextMimeTypes.TXT
                | TextMimeTypes.MD
                | TextMimeTypes.TEXT_CSV
                | TextMimeTypes.APP_CSV
            ):
                extracted_text = self.extract_from_plain_text(filepath, display_name)
            case TextMimeTypes.PDF:
                extracted_text = self.extract_from_pdf(filepath, display_name)
            case TextMimeTypes.DOCX:
                extracted_text = self.extract_from_docx(filepath, display_name)
            case TextMimeTypes.PPTX:
                extracted_text = self.extract_from_pptx(filepath, display_name)
            case TextMimeTypes.XLSX | TextMimeTypes.XLS:
                extracted_text = self.extract_from_xlsx(filepath, display_name)
            case _:
                # Fallback to plain text
                extracted_text = self.extract_from_plain_text(filepath, display_name)

        return extracted_text.strip()
